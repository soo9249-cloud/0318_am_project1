"""
badge.py — Design A / B / C
python-pptx 도형/텍스트 직접 렌더링 (Pillow 제거, 로고만 이미지 삽입)
"""
from __future__ import annotations

import io
import json
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Mm, Pt
from lxml import etree

BASE_DIR  = Path(__file__).parent
SPEC_PATH = BASE_DIR / "designs" / "design_spec.json"
with open(SPEC_PATH, encoding="utf-8") as _f:
    SPEC = json.load(_f)

A4_W, A4_H = 210.0, 297.0

BADGE_SIZES: dict[str, tuple[float, float]] = {
    "소형 70×25mm": (70.0, 25.0),
    "중형 80×25mm": (80.0, 25.0),
}
BADGE_LAYOUTS: dict[str, tuple[int, int]] = {
    "A4에 8개":  (2, 4),
    "A4에 10개": (2, 5),
    "A4에 12개": (2, 6),
}

_FONT_NAMES = {
    "kr_bold":    "굴림",
    "kr_regular": "굴림",
    "en_bold":    "Arial",
    "en_regular": "Arial",
    "en_semi":    "Arial",
}
_FONT_BOLD = {
    "kr_bold": True, "kr_regular": False,
    "en_bold": True, "en_regular": False, "en_semi": False,
}


# ─── 유틸 ─────────────────────────────────────────────────────────────────────
def _rgb(hex_color: str) -> RGBColor:
    c = hex_color.lstrip("#")
    return RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))


def _hex_t(hex_color: str) -> tuple:
    c = hex_color.lstrip("#")
    return int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)


def _darken(hex_color: str, amt: int = 40) -> str:
    r, g, b = _hex_t(hex_color)
    return "#{:02X}{:02X}{:02X}".format(max(0, r-amt), max(0, g-amt), max(0, b-amt))


def _pt_mm(pt: float) -> float:
    return pt * 25.4 / 72


def _rank_dept(person: dict, sep: str = " · ") -> str:
    rank = (person.get("rank") or "").strip()
    dept = (person.get("dept") or "").strip()
    if rank and dept:
        return f"{rank}{sep}{dept}"
    return rank or dept


# ─── 도형 헬퍼 ────────────────────────────────────────────────────────────────
def _rounded_rect(slide, x: float, y: float, w: float, h: float,
                  fill_hex: str, corner_mm: float = 1.5):
    shape = slide.shapes.add_shape(5, Mm(x), Mm(y), Mm(w), Mm(h))  # ROUNDED_RECTANGLE
    adj = min(corner_mm / (min(w, h) / 2) * 0.5, 0.5)
    shape.adjustments[0] = adj
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    shape.line.fill.background()
    return shape


def _oval(slide, x: float, y: float, d: float,
          fill_hex: str, outline_hex: Optional[str] = None, lpt: float = 0):
    shape = slide.shapes.add_shape(9, Mm(x), Mm(y), Mm(d), Mm(d))  # OVAL
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    if outline_hex and lpt:
        shape.line.color.rgb = _rgb(outline_hex)
        shape.line.width = Pt(lpt)
    else:
        shape.line.fill.background()
    return shape


def _hline(slide, x1: float, y: float, x2: float,
           hex_color: str, pt: float = 0.4):
    conn = slide.shapes.add_connector(1, Mm(x1), Mm(y), Mm(x2), Mm(y))
    conn.line.color.rgb = _rgb(hex_color)
    conn.line.width = Pt(pt)


def _vline(slide, x: float, y1: float, y2: float,
           hex_color: str, pt: float = 0.4):
    conn = slide.shapes.add_connector(1, Mm(x), Mm(y1), Mm(x), Mm(y2))
    conn.line.color.rgb = _rgb(hex_color)
    conn.line.width = Pt(pt)


# 자간 (단위: 1/100pt, 100=1pt)
_CHAR_SPACING = {
    "kr_bold":    800,
    "kr_regular":  60,
    "en_bold":      0,
    "en_regular":   0,
    "en_semi":      0,
}


def _txt(slide, x: float, y: float, w: float, h: float,
         text: str, font_key: str, size_pt: float, color_hex: str,
         align=PP_ALIGN.LEFT, all_caps: bool = False):
    """텍스트박스 하나 추가."""
    tb = slide.shapes.add_textbox(Mm(x), Mm(y), Mm(w), Mm(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = _FONT_NAMES[font_key]
    run.font.bold = _FONT_BOLD[font_key]
    run.font.size = Pt(size_pt)
    run.font.color.rgb = _rgb(color_hex)
    if all_caps:
        run.font.all_caps = True
    spc = _CHAR_SPACING.get(font_key, 0)
    if spc:
        run.font._element.set("spc", str(spc))
    return tb


_logo_cache: dict[tuple, bytes] = {}

def _logo_circle(slide, logo_path: Optional[str],
                 x: float, y: float, d: float,
                 bg_hex: str = "#FFFFFF") -> bool:
    """원형 크롭 로고 삽입 (배경색 합성). 성공 여부 반환. 결과 캐시."""
    if not logo_path:
        return False
    try:
        cache_key = ("circle", logo_path, bg_hex, d)
        if cache_key not in _logo_cache:
            from PIL import Image, ImageDraw
            px   = max(int(d * 12), 4)
            logo = Image.open(logo_path).convert("RGBA").resize((px, px))
            r, g, b = int(bg_hex[1:3], 16), int(bg_hex[3:5], 16), int(bg_hex[5:7], 16)
            bg   = Image.new("RGBA", (px, px), (r, g, b, 255))
            bg.paste(logo, mask=logo.split()[3])
            mask = Image.new("L", (px, px), 0)
            ImageDraw.Draw(mask).ellipse([0, 0, px - 1, px - 1], fill=255)
            result = Image.new("RGBA", (px, px), (0, 0, 0, 0))
            result.paste(bg, mask=mask)
            buf = io.BytesIO()
            result.save(buf, "PNG")
            _logo_cache[cache_key] = (buf.getvalue(), px, px)
        data, _, __ = _logo_cache[cache_key]
        slide.shapes.add_picture(io.BytesIO(data), Mm(x), Mm(y), Mm(d), Mm(d))
        return True
    except Exception:
        return False


def _logo_rect(slide, logo_path: Optional[str],
               x: float, y: float, w: float, h: float,
               bg_hex: str = "#FFFFFF") -> bool:
    """로고 이미지 삽입 (배경색 합성, 체커보드 방지). 원본 비율 유지, 결과 캐시."""
    if not logo_path:
        return False
    try:
        cache_key = ("rect", logo_path, bg_hex)
        if cache_key not in _logo_cache:
            from PIL import Image
            img = Image.open(logo_path).convert("RGBA")
            iw, ih = img.size
            alpha = img.split()[3]
            has_transparency = alpha.getextrema()[0] < 255

            buf = io.BytesIO()
            if has_transparency:
                img.save(buf, "PNG")
            else:
                r, g, b = int(bg_hex[1:3], 16), int(bg_hex[3:5], 16), int(bg_hex[5:7], 16)
                bg = Image.new("RGBA", img.size, (r, g, b, 255))
                bg.paste(img, mask=alpha)
                bg.convert("RGB").save(buf, "PNG")
            _logo_cache[cache_key] = (buf.getvalue(), iw, ih)
        data, iw, ih = _logo_cache[cache_key]
        aspect = iw / max(ih, 1)
        dw = min(w, h * aspect)
        dh = dw / aspect
        slide.shapes.add_picture(io.BytesIO(data),
                                 Mm(x + (w - dw) / 2), Mm(y + (h - dh) / 2),
                                 Mm(dw), Mm(dh))
        return True
    except Exception:
        return False


def _cutline(slide, x: float, y: float, w: float, h: float):
    shape = slide.shapes.add_shape(1, Mm(x), Mm(y), Mm(w), Mm(h))
    shape.fill.background()
    shape.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
    shape.line.width = Pt(0.4)
    _set_dash(shape)


def _set_dash(shape, val: str = "dash"):
    NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    ln = shape._element.find(f".//{{{NS}}}ln")
    if ln is None:
        return
    tag = f"{{{NS}}}prstDash"
    for c in list(ln):
        if c.tag == tag:
            ln.remove(c)
    etree.SubElement(ln, tag).set("val", val)


# ═══════════════════════════════════════════════════════════════════════════════
# Design A — Classic (로고 좌측 원형 + 세로 구분선 + 텍스트 우측)
# ═══════════════════════════════════════════════════════════════════════════════
def _draw_A(slide, person: dict, x: float, y: float, w: float, h: float,
            bg_hex: str, txt_hex: str, logo_path: Optional[str], company_name: str):
    PAD, GAP = 2.7, 1.2
    spec = SPEC["badge"]["A"]["text"]

    _rounded_rect(slide, x, y, w, h, bg_hex)

    left_w = w * 0.20
    div_x  = x + left_w

    # 왼쪽: 원형 로고 or 회사명 텍스트
    circle_d = min(left_w, h) - PAD * 2
    lx = x + (left_w - circle_d) / 2
    ly = y + (h - circle_d) / 2
    if not _logo_circle(slide, logo_path, lx, ly, circle_d, bg_hex=bg_hex):
        company = (company_name or "").strip()
        if company:
            _txt(slide, x + PAD * 0.3, y + PAD, left_w - PAD * 0.6, h - PAD * 2,
                 company[:5], "kr_bold", 4.5, txt_hex, PP_ALIGN.CENTER)

    # 수직 구분선
    _vline(slide, div_x, y + PAD, y + h - PAD, _darken(bg_hex, 40), 0.5)

    # 오른쪽 텍스트
    tx  = div_x + PAD
    rw  = w - (div_x - x) - PAD * 2
    sz_kr = spec["name_kr"]["size"]   # 14
    sz_en = spec["name_en"]["size"]    # 8
    sz_rd = spec["rank_dept"]["size"]  # 7

    name_kr = (person.get("name_kr") or "")
    name_en = (person.get("name_en") or "").upper()
    rd_line = _rank_dept(person)

    bh = _pt_mm(sz_kr)
    if name_en: bh += GAP + _pt_mm(sz_en)
    if rd_line: bh += GAP + 0.5 + GAP + _pt_mm(sz_rd)

    cy = y + (h - bh) / 2

    if name_kr:
        _txt(slide, tx, cy, rw, _pt_mm(sz_kr) + 1, name_kr, "kr_bold", sz_kr, txt_hex)
        cy += _pt_mm(sz_kr) + GAP
    if name_en:
        _txt(slide, tx, cy, rw, _pt_mm(sz_en) + 1, name_en,
             "en_bold", sz_en, txt_hex, all_caps=True)
        cy += _pt_mm(sz_en) + GAP
    if rd_line:
        muted = _darken(txt_hex, 30)
        _hline(slide, tx, cy + 0.25, tx + rw * 0.85, muted, 0.3)
        cy += 0.5 + GAP
        _txt(slide, tx, cy, rw, _pt_mm(sz_rd) + 1, rd_line,
             "kr_regular", sz_rd, muted)


# ═══════════════════════════════════════════════════════════════════════════════
# Design B — Centered (전체 중앙 + 가로 구분선)
# ═══════════════════════════════════════════════════════════════════════════════
def _draw_B(slide, person: dict, x: float, y: float, w: float, h: float,
            bg_hex: str, txt_hex: str, logo_path: Optional[str], company_name: str):
    PAD, GAP = 2.25, 1.0
    spec   = SPEC["badge"]["B"]["text"]
    margin = w * 0.15
    rw     = w - margin * 2

    _rounded_rect(slide, x, y, w, h, bg_hex)

    sz_kr = spec["name_kr"]["size"]  # 16
    sz_en = spec["name_en"]["size"]   # 9
    sz_rk = spec["rank"]["size"]       # 8
    sz_co = spec["company"]["size"]    # 7

    name_kr = (person.get("name_kr") or "")
    name_en = (person.get("name_en") or "").upper()
    rank    = _rank_dept(person, " · ")
    company = (person.get("company") or company_name or "").upper()

    top_lines = []
    if name_kr: top_lines.append((name_kr, "kr_bold",    sz_kr, False))
    if name_en: top_lines.append((name_en, "en_bold",    sz_en, True))
    if rank:    top_lines.append((rank,    "kr_regular",  sz_rk, False))

    top_h    = sum(_pt_mm(sz) + GAP for _, _, sz, _ in top_lines) - (GAP if top_lines else 0)
    bot_h    = _pt_mm(sz_co)
    rule_gap = GAP * 2
    total_h  = top_h + rule_gap + 0.2 + rule_gap + bot_h

    cy = y + (h - total_h) / 2

    for text, fk, sz, caps in top_lines:
        _txt(slide, x + margin, cy, rw, _pt_mm(sz) + 1,
             text, fk, sz, txt_hex, PP_ALIGN.CENTER, all_caps=caps)
        cy += _pt_mm(sz) + GAP

    cy += rule_gap - GAP
    _hline(slide, x + margin, cy, x + w - margin, _darken(txt_hex, 30), 0.3)
    cy += 0.2 + rule_gap

    muted = _darken(txt_hex, 30)
    if logo_path:
        logo_h = bot_h * 2.5
        logo_w = min(rw, logo_h * 4)
        _logo_rect(slide, logo_path, x + margin + (rw - logo_w) / 2, cy, logo_w, logo_h, bg_hex=bg_hex)
    elif company:
        _txt(slide, x + margin, cy, rw, _pt_mm(sz_co) + 1,
             company, "en_semi", sz_co, muted, PP_ALIGN.CENTER, all_caps=True)


# ═══════════════════════════════════════════════════════════════════════════════
# Design C — Modern (미니멀 원형 로고 + 우측 텍스트, 여백 넉넉)
# ═══════════════════════════════════════════════════════════════════════════════
def _draw_C(slide, person: dict, x: float, y: float, w: float, h: float,
            bg_hex: str, txt_hex: str, logo_path: Optional[str], company_name: str):
    GPAD, GAP = 3.75, 1.3
    spec = SPEC["badge"]["C"]["text"]

    _rounded_rect(slide, x, y, w, h, bg_hex)

    sz_kr = spec["name_kr"]["size"]  # 15
    sz_rk = spec["rank"]["size"]       # 8
    sz_dp = spec["dept"]["size"]       # 7

    name_kr  = (person.get("name_kr") or "")
    rank     = (person.get("rank") or "").strip()
    dept     = (person.get("dept") or "").strip()
    rank_ln  = rank if rank else dept
    dept_ln  = dept if (rank and dept) else ""
    muted    = _darken(txt_hex, 40)

    # 원형 로고 or 색상 점
    circle_d = h - GPAD * 2
    lx = x + GPAD
    ly = y + GPAD

    if not _logo_circle(slide, logo_path, lx, ly, circle_d, bg_hex=bg_hex):
        # Fallback: 배경보다 약간 밝은 원
        r, g, b = _hex_t(bg_hex)
        dot = "#{:02X}{:02X}{:02X}".format(
            min(255, r + 40), min(255, g + 40), min(255, b + 40))
        _oval(slide, lx, ly, circle_d, dot)

    tx = lx + circle_d + GPAD
    rw = w - (tx - x) - GPAD

    lines = []
    if name_kr:  lines.append((name_kr,  "kr_bold",    sz_kr, txt_hex))
    if rank_ln:  lines.append((rank_ln,  "kr_regular", sz_rk, txt_hex))
    if dept_ln:  lines.append((dept_ln,  "kr_regular", sz_dp, muted))

    bh = sum(_pt_mm(sz) + GAP for _, _, sz, _ in lines) - (GAP if lines else 0)
    cy = y + (h - bh) / 2

    for text, fk, sz, col in lines:
        _txt(slide, tx, cy, rw, _pt_mm(sz) + 1, text, fk, sz, col)
        cy += _pt_mm(sz) + GAP


# ═══════════════════════════════════════════════════════════════════════════════
# 공개 API
# ═══════════════════════════════════════════════════════════════════════════════
_DRAW = {"A": _draw_A, "B": _draw_B, "C": _draw_C}


def place_on_slide(
    slide,
    people:        list[dict],
    layout:        str,
    badge_size_mm: tuple[float, float],
    design:        str,
    color:         str,
    logo_path:     Optional[str],
    company_name:  str,
) -> None:
    """A4 슬라이드에 배지 배치 (python-pptx 도형/텍스트 직접 렌더링)."""
    cols, rows = BADGE_LAYOUTS.get(layout, (2, 4))
    bw, bh     = badge_size_mm

    GAP_X = 5.0  # mm
    GAP_Y = 8.0  # mm

    margin_x = (A4_W - cols * bw - (cols - 1) * GAP_X) / 2
    margin_y = (A4_H - rows * bh - (rows - 1) * GAP_Y) / 2

    colors  = SPEC["badge"][design]["colors"][color]
    bg_hex  = colors["bg"]
    txt_hex = colors["text"]
    draw_fn = _DRAW[design]

    for idx, person in enumerate(people[: cols * rows]):
        row, col = divmod(idx, cols)
        lx = margin_x + col * (bw + GAP_X)
        ty = margin_y + row * (bh + GAP_Y)
        draw_fn(slide, person, lx, ty, bw, bh, bg_hex, txt_hex, logo_path, company_name)
        _cutline(slide, lx, ty, bw, bh)


# ═══════════════════════════════════════════════════════════════════════════════
# 엑셀 읽기
# ═══════════════════════════════════════════════════════════════════════════════
def read_excel(path: str) -> list[dict]:
    """엑셀 파일을 읽어 사람 목록 반환. 실패 시 [{"error": "..."}]."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True)
        ws = wb.active

        headers = [str(cell.value or "").strip() for cell in ws[1]]
        KEY_MAP = {
            "이름": "name_kr", "한글이름": "name_kr", "name_kr": "name_kr",
            "영문이름": "name_en", "영어이름": "name_en", "영문": "name_en", "name_en": "name_en",
            "직급": "rank", "직위": "rank", "rank": "rank",
            "부서": "dept", "부서명": "dept", "department": "dept", "dept": "dept",
            "회사": "company", "회사명": "company", "company": "company",
        }
        col_keys = [KEY_MAP.get(h.lower(), h.lower()) for h in headers]

        people = []
        for row in ws.iter_rows(min_row=2, values_only=True):
            person = {col_keys[i]: str(v or "").strip()
                      for i, v in enumerate(row) if i < len(col_keys)}
            if any(person.values()):
                people.append(person)

        if not people:
            return [{"error": "데이터가 없습니다. 엑셀 파일을 확인해주세요."}]
        return people
    except Exception as e:
        return [{"error": f"엑셀 읽기 실패: {e}"}]


# ═══════════════════════════════════════════════════════════════════════════════
# PPTX 생성
# ═══════════════════════════════════════════════════════════════════════════════
def generate_badges(
    people:       list[dict],
    logo_path:    Optional[str],
    company_name: str,
    badge_type:   str,
    design:       str,
    color:        str,
    size:         str,
    layout:       str,
    event_mode:   bool,
    event_info:   Optional[dict],
    output_dir:   str,
) -> Optional[str]:
    """PPTX 생성 후 파일 경로 반환. 실패 시 None."""
    import traceback
    from datetime import datetime

    try:
        prs = Presentation()
        blank = prs.slide_layouts[6]
        is_lanyard = badge_type == "목걸이형" or design in ("D", "E", "F")

        if is_lanyard:
            from designs import lanyard as ln

            badge_size   = ln.LANYARD_SIZES.get(size, (90.0, 120.0))
            is_landscape = (size == "대형 103×133mm")

            if is_landscape:
                slide_w, slide_h = 297.0, 210.0
                layout_key = "대형_2개"
            else:
                slide_w, slide_h = 210.0, 297.0
                _layout_map = {"A4에 2개": "일반_2개", "A4에 4개": "일반_4개"}
                layout_key = _layout_map.get(layout, layout if layout in ln.LANYARD_LAYOUTS else "일반_2개")

            prs.slide_width  = Mm(slide_w)
            prs.slide_height = Mm(slide_h)

            cols, rows = ln.LANYARD_LAYOUTS[layout_key]
            per_page   = cols * rows

            for start in range(0, len(people), per_page):
                batch = people[start: start + per_page]
                slide = prs.slides.add_slide(blank)
                ln.place_on_slide(
                    slide, batch, layout_key, badge_size,
                    design, color, logo_path, company_name,
                    event_mode, event_info,
                    slide_w=slide_w, slide_h=slide_h,
                )
        else:
            badge_size = BADGE_SIZES.get(size, (80.0, 25.0))
            prs.slide_width  = Mm(A4_W)
            prs.slide_height = Mm(A4_H)

            cols, rows = BADGE_LAYOUTS.get(layout, (2, 4))
            per_page   = cols * rows

            for start in range(0, len(people), per_page):
                batch = people[start: start + per_page]
                slide = prs.slides.add_slide(blank)
                place_on_slide(slide, batch, layout, badge_size,
                               design, color, logo_path, company_name)

        ts       = datetime.now().strftime("%Y%m%d_%H%M%S")
        out_path = str(Path(output_dir) / f"badges_{ts}.pptx")
        prs.save(out_path)

        try:
            from font_embed import embed_fonts
            embed_fonts(out_path, BASE_DIR / "fonts")
        except Exception:
            pass

        return out_path
    except Exception:
        traceback.print_exc()
        return None

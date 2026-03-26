"""
lanyard.py — Design D / E / F
python-pptx 도형/텍스트 직접 렌더링 (로고만 이미지 삽입)
"""
from __future__ import annotations

import io
import json
from pathlib import Path
from typing import Optional

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Mm, Pt
from lxml import etree

BASE_DIR  = Path(__file__).parent.parent
SPEC_PATH = BASE_DIR / "designs" / "design_spec.json"
with open(SPEC_PATH, encoding="utf-8") as _f:
    SPEC = json.load(_f)

A4_W, A4_H = 210.0, 297.0

LANYARD_SIZES: dict[str, tuple[float, float]] = {
    "일반 90×120mm":  (90.0,  120.0),
    "대형 103×133mm": (103.0, 133.0),
}
LANYARD_LAYOUTS: dict[str, tuple[int, int]] = {
    "일반_2개": (1, 2),   # 세로 A4, 1열 2행
    "일반_4개": (2, 2),   # 세로 A4, 2열 2행
    "대형_2개": (2, 1),   # 가로 A4, 2열 1행
}
BADGE_GAP = 10.0  # 명찰 간 여백 (mm)

_FONT_NAMES = {
    "kr_bold":    "굴림",
    "kr_regular": "굴림",
    "en_bold":    "Arial",
    "en_regular": "Arial",
    "en_semi":    "Arial",
}
_FONT_BOLD = {
    "kr_bold": True, "kr_regular": False,
    "en_bold": True, "en_regular": False, "en_semi": False,
}


# ─── 유틸 ─────────────────────────────────────────────────────────────────────
def _rgb(hex_color: str) -> RGBColor:
    c = hex_color.lstrip("#")[:6]
    return RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))


def _pt_mm(pt: float) -> float:
    return pt * 25.4 / 72


def _rank_dept(person: dict, sep: str = " · ") -> str:
    rank = (person.get("rank") or "").strip()
    dept = (person.get("dept") or "").strip()
    if rank and dept:
        return f"{rank}{sep}{dept}"
    return rank or dept


# ─── 도형 헬퍼 ────────────────────────────────────────────────────────────────
def _rounded_rect(slide, x: float, y: float, w: float, h: float,
                  fill_hex: str, corner_mm: float = 2.0):
    shape = slide.shapes.add_shape(5, Mm(x), Mm(y), Mm(w), Mm(h))
    adj = min(corner_mm / (min(w, h) / 2) * 0.5, 0.5)
    shape.adjustments[0] = adj
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    shape.line.fill.background()
    return shape


def _rect(slide, x: float, y: float, w: float, h: float, fill_hex: str):
    shape = slide.shapes.add_shape(1, Mm(x), Mm(y), Mm(w), Mm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    shape.line.fill.background()
    return shape


def _oval(slide, x: float, y: float, w: float, h: float,
          fill_hex: str, outline_hex: Optional[str] = None, lpt: float = 0):
    shape = slide.shapes.add_shape(9, Mm(x), Mm(y), Mm(w), Mm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    if outline_hex and lpt:
        shape.line.color.rgb = _rgb(outline_hex)
        shape.line.width = Pt(lpt)
    else:
        shape.line.fill.background()
    return shape


def _set_gradient(shape, hex1: str, hex2: str, angle_60k: int = 5400000):
    """수직 그라데이션 fill (XML 직접 조작)."""
    spPr = shape._element.spPr
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for child in list(spPr):
        tag = child.tag.split("}")[-1]
        if tag in ("solidFill", "gradFill", "noFill", "blipFill", "pattFill"):
            spPr.remove(child)
    gf  = etree.SubElement(spPr, f"{{{A}}}gradFill")
    gsl = etree.SubElement(gf,  f"{{{A}}}gsLst")
    gs1 = etree.SubElement(gsl, f"{{{A}}}gs"); gs1.set("pos", "0")
    etree.SubElement(gs1, f"{{{A}}}srgbClr").set("val", hex1.lstrip("#"))
    gs2 = etree.SubElement(gsl, f"{{{A}}}gs"); gs2.set("pos", "100000")
    etree.SubElement(gs2, f"{{{A}}}srgbClr").set("val", hex2.lstrip("#"))
    lin = etree.SubElement(gf, f"{{{A}}}lin")
    lin.set("ang", str(angle_60k))
    lin.set("scaled", "0")
    shape.line.fill.background()


def _hline(slide, x1: float, y: float, x2: float,
           hex_color: str, pt: float = 0.4):
    conn = slide.shapes.add_connector(1, Mm(x1), Mm(y), Mm(x2), Mm(y))
    conn.line.color.rgb = _rgb(hex_color)
    conn.line.width = Pt(pt)


# 자간 (단위: 1/100pt, 100=1pt)
_CHAR_SPACING = {
    "kr_bold":    800,
    "kr_regular":  60,
    "en_bold":      0,
    "en_regular":   0,
    "en_semi":      0,
}


def _txt(slide, x: float, y: float, w: float, h: float,
         text: str, font_key: str, size_pt: float, color_hex: str,
         align=PP_ALIGN.CENTER, all_caps: bool = False):
    tb = slide.shapes.add_textbox(Mm(x), Mm(y), Mm(w), Mm(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = _FONT_NAMES[font_key]
    run.font.bold = _FONT_BOLD[font_key]
    run.font.size = Pt(size_pt)
    run.font.color.rgb = _rgb(color_hex[:7])
    if all_caps:
        run.font.all_caps = True
    spc = _CHAR_SPACING.get(font_key, 0)
    if spc:
        run.font._element.set("spc", str(spc))
    return tb


_logo_cache: dict[tuple, bytes] = {}

def _logo_rect(slide, logo_path: Optional[str],
               x: float, y: float, w: float, h: float,
               bg_hex: str = "#FFFFFF") -> bool:
    """로고를 배경색 위에 합성해서 삽입 (투명도 체커보드 방지). 결과 캐시."""
    if not logo_path:
        return False
    try:
        cache_key = (logo_path, bg_hex)
        if cache_key not in _logo_cache:
            from PIL import Image
            img = Image.open(logo_path).convert("RGBA")
            r, g, b = int(bg_hex[1:3], 16), int(bg_hex[3:5], 16), int(bg_hex[5:7], 16)
            bg = Image.new("RGBA", img.size, (r, g, b, 255))
            bg.paste(img, mask=img.split()[3])
            buf = io.BytesIO()
            bg.convert("RGB").save(buf, "PNG")
            _logo_cache[cache_key] = buf.getvalue()
        slide.shapes.add_picture(io.BytesIO(_logo_cache[cache_key]), Mm(x), Mm(y), Mm(w), Mm(h))
        return True
    except Exception:
        return False


def _cutline(slide, x: float, y: float, w: float, h: float):
    shape = slide.shapes.add_shape(1, Mm(x), Mm(y), Mm(w), Mm(h))
    shape.fill.background()
    shape.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
    shape.line.width = Pt(0.4)
    _set_dash(shape)


def _set_dash(shape, val: str = "dash"):
    NS  = "http://schemas.openxmlformats.org/drawingml/2006/main"
    ln  = shape._element.find(f".//{{{NS}}}ln")
    if ln is None:
        return
    tag = f"{{{NS}}}prstDash"
    for c in list(ln):
        if c.tag == tag:
            ln.remove(c)
    etree.SubElement(ln, tag).set("val", val)


def _hole(slide, badge_x: float, badge_y: float, badge_w: float,
          bg_is_dark: bool = False):
    """상단 중앙 구멍 표시."""
    r    = SPEC["common"]["hole"]["radius_mm"]
    cx   = badge_x + badge_w / 2
    cy   = badge_y + 2.5 + r
    fill = "#B4B4B4" if not bg_is_dark else "#DDDDDD"
    outl = "#5A5A5A" if not bg_is_dark else "#AAAAAA"
    _oval(slide, cx - r, cy - r, r * 2, r * 2, fill, outl, 1.0)


def _draw_wave_band(slide, x: float, wave_start_y: float,
                    w: float, band_h: float, fill_hex: str):
    """
    웨이브 모양 하단 밴드.
    컬러 직사각형 위에 흰색 타원을 겹쳐서 물결 효과 구현.
    """
    # 1. 컬러 밴드 (직사각형)
    _rect(slide, x, wave_start_y, w, band_h, fill_hex)

    # 2. 흰색 돔(타원)을 위에 겹쳐서 웨이브 상단 곡선 표현
    dome_h = band_h * 0.55     # 타원 높이 (밴드 높이의 55%)
    dome_w = w * 1.05          # 타원 너비 (카드보다 약간 넓게)
    dome_x = x - (dome_w - w) / 2
    dome_y = wave_start_y - dome_h * 0.5   # 절반쯤 위로 올림

    shape = slide.shapes.add_shape(9, Mm(dome_x), Mm(dome_y), Mm(dome_w), Mm(dome_h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb("#FFFFFF")
    shape.line.fill.background()


# ═══════════════════════════════════════════════════════════════════════════════
# Design D — Wave Corporate (흰 배경 + 하단 웨이브 컬러 밴드)
# ═══════════════════════════════════════════════════════════════════════════════
def _draw_D(slide, person: dict, x: float, y: float, w: float, h: float,
            wave_hex: str, logo_path: Optional[str], company_name: str,
            event_mode: bool, event_info: Optional[dict]):
    PAD, GAP = 5.0, 2.0
    hole_offset = 12.0
    spec = SPEC["lanyard"]["D"]["text"]

    _rounded_rect(slide, x, y, w, h, "#FFFFFF")

    bot_h = h * 0.25
    bot_y = y + h - bot_h

    # ── 하단 웨이브 밴드 (항상)
    _draw_wave_band(slide, x, bot_y, w, bot_h, wave_hex)

    # ── 하단 밴드: 회사명 (항상)
    sz_co   = spec["company"]["size"]
    company = (person.get("company") or company_name or "").upper()
    if company:
        cy_co = bot_y + (bot_h - _pt_mm(sz_co)) / 2 + bot_h * 0.1
        _txt(slide, x + PAD, cy_co, w - PAD * 2, _pt_mm(sz_co) + 1,
             company, "en_semi", sz_co, "#FFFFFF", all_caps=True)

    sz_kr = spec["name_kr"]["size"]
    sz_en = spec["name_en"]["size"]
    sz_dr = spec["dept_rank"]["size"]
    name_kr = (person.get("name_kr") or "")
    name_en = (person.get("name_en") or "").upper()
    dept_rk = _rank_dept(person)

    if event_mode and event_info:
        # ── 행사 모드: 상단(로고+행사정보) / 중단(이름) 구조
        top_h = h * 0.28
        mid_h = bot_y - (y + top_h)
        mid_y = y + top_h

        _hline(slide, x + PAD * 2, mid_y, x + w - PAD * 2, "#E0E0E0", 0.4)

        sz_ev   = SPEC["common"]["event_mode"]["font_size"]
        ev_name = (event_info.get("event_name") or "").strip()
        ev_date = (event_info.get("event_date") or "").strip()

        top_cy = y + hole_offset + (top_h - hole_offset) / 2
        if logo_path:
            logo_h = min((top_h - hole_offset) * 0.42, 13.0)
            logo_w = min(w * 0.50, logo_h * 4.0)
            logo_y = y + hole_offset + 1.0
            _logo_rect(slide, logo_path,
                       x + (w - logo_w) / 2, logo_y, logo_w, logo_h,
                       bg_hex="#FFFFFF")
            ey = logo_y + logo_h + 1.5
        elif company_name:
            _txt(slide, x + PAD, top_cy - _pt_mm(11) / 2, w - PAD * 2,
                 _pt_mm(11) + 1, company_name, "kr_bold", 11, "#1E1E1E")
            ey = top_cy + _pt_mm(11) / 2 + 1.5
        else:
            ey = top_cy - (_pt_mm(sz_ev) + (GAP + _pt_mm(sz_ev * 0.75) if ev_date else 0)) / 2

        if ev_name:
            _txt(slide, x + PAD, ey, w - PAD * 2, _pt_mm(sz_ev) + 1,
                 ev_name, "kr_bold", sz_ev, "#1E1E1E")
            ey += _pt_mm(sz_ev) + 1.5
        if ev_date:
            _txt(slide, x + PAD, ey, w - PAD * 2, _pt_mm(sz_ev * 0.75) + 1,
                 ev_date, "kr_bold", sz_ev * 0.75, "#505050")

        lines = []
        if name_kr: lines.append((name_kr, "kr_bold",    sz_kr, "#141414", False))
        if name_en: lines.append((name_en, "en_bold",    sz_en, "#3C3C3C", True))
        if dept_rk: lines.append((dept_rk, "kr_regular", sz_dr, "#505050", False))
        bh_txt = sum(_pt_mm(sz) + GAP for _, _, sz, _, _ in lines) - (GAP if lines else 0)
        cy = mid_y + (mid_h - bh_txt) / 2
        for text, fk, sz, col, caps in lines:
            _txt(slide, x + PAD, cy, w - PAD * 2, _pt_mm(sz) + 1,
                 text, fk, sz, col, all_caps=caps)
            cy += _pt_mm(sz) + GAP

    else:
        # ── 일반 모드: 로고(중앙 살짝 위) → 이름 → 부서
        logo_h   = min(h * 0.20, 22.0)
        logo_w   = min(w * 0.60, logo_h * 4.0)
        logo_gap = 4.0

        lines = []
        if name_kr: lines.append((name_kr, "kr_bold",    sz_kr, "#141414", False))
        if name_en: lines.append((name_en, "en_bold",    sz_en, "#3C3C3C", True))
        if dept_rk: lines.append((dept_rk, "kr_regular", sz_dr, "#505050", False))

        logo_cy  = y + h * 0.42
        logo_top = max(logo_cy - logo_h / 2, y + hole_offset + 2.0)
        if logo_path:
            _logo_rect(slide, logo_path,
                       x + (w - logo_w) / 2, logo_top, logo_w, logo_h,
                       bg_hex="#FFFFFF")
            cy = logo_top + logo_h + logo_gap
        elif company_name:
            _txt(slide, x + PAD, logo_top + (logo_h - _pt_mm(13)) / 2, w - PAD * 2,
                 _pt_mm(13) + 1, company_name, "kr_bold", 13, "#1E1E1E")
            cy = logo_top + logo_h + logo_gap
        else:
            cy = y + h * 0.44

        for text, fk, sz, col, caps in lines:
            _txt(slide, x + PAD, cy, w - PAD * 2, _pt_mm(sz) + 1,
                 text, fk, sz, col, all_caps=caps)
            cy += _pt_mm(sz) + GAP

    _hole(slide, x, y, w, bg_is_dark=False)


# ═══════════════════════════════════════════════════════════════════════════════
# Design E — Gradient Vivid (그라데이션 배경 + 흰 텍스트)
# ═══════════════════════════════════════════════════════════════════════════════
def _draw_E(slide, person: dict, x: float, y: float, w: float, h: float,
            grad_colors: list[str], logo_path: Optional[str], company_name: str,
            event_mode: bool, event_info: Optional[dict]):
    PAD, GAP = 5.0, 2.5
    hole_offset = 14.0
    spec = SPEC["lanyard"]["E"]["text"]

    # 그라데이션 배경
    bg_shape = slide.shapes.add_shape(5, Mm(x), Mm(y), Mm(w), Mm(h))
    adj = min(2.0 / (min(w, h) / 2) * 0.5, 0.5)
    bg_shape.adjustments[0] = adj
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(grad_colors[0])
    bg_shape.line.fill.background()
    _set_gradient(bg_shape, grad_colors[0], grad_colors[1])

    sz_kr = spec["name_kr"]["size"]
    sz_en = spec["name_en"]["size"]
    sz_dr = spec["dept_rank"]["size"]
    name_kr = (person.get("name_kr") or "")
    name_en = (person.get("name_en") or "").upper()
    dept_rk = _rank_dept(person)

    if event_mode and event_info:
        # ── 행사 모드: 상단(로고/회사) + 하단(행사정보+이름) 기존 구조 유지
        top_area_h = h * 0.25
        top_cy     = y + hole_offset + (top_area_h - hole_offset) / 2

        if logo_path:
            logo_h = top_area_h * 0.5
            logo_w = min(w * 0.5, logo_h * 4)
            _logo_rect(slide, logo_path,
                       x + (w - logo_w) / 2, top_cy - logo_h / 2,
                       logo_w, logo_h, bg_hex=grad_colors[0])
        else:
            company = (company_name or "").strip()
            if company:
                _txt(slide, x + PAD, top_cy - _pt_mm(11) / 2, w - PAD * 2,
                     _pt_mm(11) + 1, company, "en_semi", 11, "#EEEEEE")

        lines = []
        sz_ev   = SPEC["common"]["event_mode"]["font_size"]
        ev_name = (event_info.get("event_name") or "").strip()
        ev_date = (event_info.get("event_date") or "").strip()
        if ev_name: lines.append((ev_name, "kr_bold", sz_ev, "#FFFFFF", False))
        if ev_date: lines.append((ev_date, "en_semi", 10,    "#EEEEEE", False))
        if name_kr: lines.append((name_kr, "kr_bold",    sz_kr, "#FFFFFF", False))
        if name_en: lines.append((name_en, "en_bold",    sz_en, "#EEEEEE", True))
        if dept_rk: lines.append((dept_rk, "kr_regular", sz_dr, "#EEEEEE", False))

        bh = sum(_pt_mm(sz) + GAP for _, _, sz, _, _ in lines) - (GAP if lines else 0)
        usable_top = y + top_area_h
        usable_h   = h - top_area_h - h * 0.08
        cy = usable_top + (usable_h - bh) / 2
        for text, fk, sz, col, caps in lines:
            _txt(slide, x + PAD, cy, w - PAD * 2, _pt_mm(sz) + 1,
                 text, fk, sz, col, all_caps=caps)
            cy += _pt_mm(sz) + GAP

    else:
        # ── 일반 모드: 로고(중앙 살짝 위) → 이름 → 부서
        logo_h   = min(h * 0.20, 22.0)
        logo_w   = min(w * 0.60, logo_h * 4.0)
        logo_gap = 4.0

        lines = []
        if name_kr: lines.append((name_kr, "kr_bold",    sz_kr, "#FFFFFF", False))
        if name_en: lines.append((name_en, "en_bold",    sz_en, "#EEEEEE", True))
        if dept_rk: lines.append((dept_rk, "kr_regular", sz_dr, "#EEEEEE", False))

        logo_cy  = y + h * 0.42
        logo_top = max(logo_cy - logo_h / 2, y + hole_offset + 2.0)
        if logo_path:
            _logo_rect(slide, logo_path,
                       x + (w - logo_w) / 2, logo_top, logo_w, logo_h,
                       bg_hex=grad_colors[0])
            cy = logo_top + logo_h + logo_gap
        elif company_name:
            _txt(slide, x + PAD, logo_top + (logo_h - _pt_mm(13)) / 2, w - PAD * 2,
                 _pt_mm(13) + 1, company_name, "kr_bold", 13, "#FFFFFF")
            cy = logo_top + logo_h + logo_gap
        else:
            cy = y + h * 0.44

        for text, fk, sz, col, caps in lines:
            _txt(slide, x + PAD, cy, w - PAD * 2, _pt_mm(sz) + 1,
                 text, fk, sz, col, all_caps=caps)
            cy += _pt_mm(sz) + GAP

    _hole(slide, x, y, w, bg_is_dark=True)


# ═══════════════════════════════════════════════════════════════════════════════
# Design F — Bold Solid (단색 + 흰 상단 바)
# ═══════════════════════════════════════════════════════════════════════════════
def _draw_F(slide, person: dict, x: float, y: float, w: float, h: float,
            bg_hex: str, logo_path: Optional[str], company_name: str,
            event_mode: bool, event_info: Optional[dict]):
    PAD, GAP = 5.0, 2.5
    spec = SPEC["lanyard"]["F"]["text"]

    _rounded_rect(slide, x, y, w, h, bg_hex)

    sz_kr = spec["name_kr"]["size"]
    sz_en = spec["name_en"]["size"]
    sz_dr = spec["dept_rank"]["size"]
    name_kr = (person.get("name_kr") or "")
    name_en = (person.get("name_en") or "").upper()
    dept_rk = _rank_dept(person)

    if event_mode and event_info:
        # ── 행사 모드: 흰 상단 바(로고/회사명+행사정보) + 중단(이름) + 하단(부서)
        top_h = h * 0.30
        mid_h = h * 0.45
        bot_h = h - top_h - mid_h
        mid_y = y + top_h
        bot_y = mid_y + mid_h
        hole_offset = 12.0

        _rect(slide, x, y, w, top_h, "#FFFFFF")
        _hline(slide, x, mid_y, x + w, "#DCDCDC", 0.3)

        sz_ev   = SPEC["common"]["event_mode"]["font_size"] * 0.85
        ev_name = (event_info.get("event_name") or "").strip()
        ev_date = (event_info.get("event_date") or "").strip()

        if logo_path:
            logo_h = min((top_h - hole_offset) * 0.40, 12.0)
            logo_w = min(w * 0.50, logo_h * 4.0)
            logo_y = y + hole_offset + 1.0
            _logo_rect(slide, logo_path,
                       x + (w - logo_w) / 2, logo_y, logo_w, logo_h,
                       bg_hex="#FFFFFF")
            ey = logo_y + logo_h + 1.0
        elif company_name:
            co_y = y + hole_offset + 1.0
            _txt(slide, x + PAD, co_y, w - PAD * 2, _pt_mm(10) + 1,
                 company_name, "kr_bold", 10, bg_hex)
            ey = co_y + _pt_mm(10) + 1.5
        else:
            ey = y + hole_offset + (top_h - hole_offset) / 2 - _pt_mm(sz_ev) / 2

        if ev_name:
            _txt(slide, x + PAD, ey, w - PAD * 2, _pt_mm(sz_ev) + 1,
                 ev_name, "kr_bold", sz_ev, bg_hex)
            ey += _pt_mm(sz_ev) + 1.0
        if ev_date:
            _txt(slide, x + PAD, ey, w - PAD * 2, _pt_mm(9) + 1,
                 ev_date, "en_semi", 9, "#646464")

        mid_lines = []
        if name_kr: mid_lines.append((name_kr, "kr_bold", sz_kr, False))
        if name_en: mid_lines.append((name_en, "en_bold", sz_en, True))
        bh = sum(_pt_mm(sz) + GAP for _, _, sz, _ in mid_lines) - (GAP if mid_lines else 0)
        cy = mid_y + (mid_h - bh) / 2
        for text, fk, sz, caps in mid_lines:
            _txt(slide, x + PAD, cy, w - PAD * 2, _pt_mm(sz) + 1,
                 text, fk, sz, "#FFFFFF", all_caps=caps)
            cy += _pt_mm(sz) + GAP

        if dept_rk:
            cy_bot = bot_y + (bot_h - _pt_mm(sz_dr)) / 2
            _txt(slide, x + PAD, cy_bot, w - PAD * 2, _pt_mm(sz_dr) + 1,
                 dept_rk, "kr_regular", sz_dr, "#EEEEEE")

    else:
        # ── 일반 모드: 흰 상단 바에 로고, 로고(중앙 살짝 위) → 이름 → 부서
        hole_offset = 12.0
        logo_h   = min(h * 0.20, 22.0)
        logo_w   = min(w * 0.60, logo_h * 4.0)
        logo_gap = 4.0

        # 흰 상단 바: 로고 영역에 대비를 줌
        white_bar_h = h * 0.45
        _rect(slide, x, y, w, white_bar_h, "#FFFFFF")
        _hline(slide, x, y + white_bar_h, x + w, "#DCDCDC", 0.3)

        lines = []
        if name_kr: lines.append((name_kr, "kr_bold",    sz_kr, "#FFFFFF", False))
        if name_en: lines.append((name_en, "en_bold",    sz_en, "#EEEEEE", True))
        if dept_rk: lines.append((dept_rk, "kr_regular", sz_dr, "#EEEEEE", False))

        logo_cy  = y + h * 0.42
        logo_top = max(logo_cy - logo_h / 2, y + hole_offset + 2.0)
        if logo_path:
            _logo_rect(slide, logo_path,
                       x + (w - logo_w) / 2, logo_top, logo_w, logo_h,
                       bg_hex="#FFFFFF")
            cy = logo_top + logo_h + logo_gap
        elif company_name:
            _txt(slide, x + PAD, logo_top + (logo_h - _pt_mm(13)) / 2, w - PAD * 2,
                 _pt_mm(13) + 1, company_name, "kr_bold", 13, bg_hex)
            cy = logo_top + logo_h + logo_gap
        else:
            cy = y + h * 0.44

        for text, fk, sz, col, caps in lines:
            # 흰 바 위: bg_hex 색상, 흰 바 아래: 흰색
            text_col = bg_hex if (cy + _pt_mm(sz) / 2) < (y + white_bar_h) else col
            _txt(slide, x + PAD, cy, w - PAD * 2, _pt_mm(sz) + 1,
                 text, fk, sz, text_col, all_caps=caps)
            cy += _pt_mm(sz) + GAP

    _hole(slide, x, y, w, bg_is_dark=True)


# ═══════════════════════════════════════════════════════════════════════════════
# 공개 API
# ═══════════════════════════════════════════════════════════════════════════════
def place_on_slide(
    slide,
    people:        list[dict],
    layout:        str,
    badge_size_mm: tuple[float, float],
    design:        str,
    color:         str,
    logo_path:     Optional[str],
    company_name:  str,
    event_mode:    bool = False,
    event_info:    Optional[dict] = None,
    slide_w:       float = A4_W,
    slide_h:       float = A4_H,
) -> None:
    """슬라이드에 랜야드 배치."""
    cols, rows = LANYARD_LAYOUTS.get(layout, (1, 2))
    bw, bh     = badge_size_mm
    gap        = BADGE_GAP

    margin_x = (slide_w - cols * bw - (cols - 1) * gap) / 2
    margin_y = (slide_h - rows * bh - (rows - 1) * gap) / 2

    for idx, person in enumerate(people[: cols * rows]):
        row, col = divmod(idx, cols)
        lx = margin_x + col * (bw + gap)
        ty = margin_y + row * (bh + gap)

        if design == "D":
            wave_colors = SPEC["lanyard"]["D"]["wave_colors"]
            wave_hex    = wave_colors.get(color, next(iter(wave_colors.values())))
            _draw_D(slide, person, lx, ty, bw, bh, wave_hex,
                    logo_path, company_name, event_mode, event_info)
        elif design == "E":
            gradients = SPEC["lanyard"]["E"]["gradients"]
            grad_pair = gradients.get(color, next(iter(gradients.values())))
            _draw_E(slide, person, lx, ty, bw, bh, grad_pair,
                    logo_path, company_name, event_mode, event_info)
        elif design == "F":
            bg_colors = SPEC["lanyard"]["F"]["bg_colors"]
            bg_hex    = bg_colors.get(color, next(iter(bg_colors.values())))
            _draw_F(slide, person, lx, ty, bw, bh, bg_hex,
                    logo_path, company_name, event_mode, event_info)

        _cutline(slide, lx, ty, bw, bh)
